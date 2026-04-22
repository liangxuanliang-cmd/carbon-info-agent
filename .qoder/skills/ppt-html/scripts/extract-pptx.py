#!/usr/bin/env python3
"""
从 PowerPoint 文件（.pptx）中提取所有内容。
返回包含幻灯片、文本和图片的 JSON 结构。

用法：
    python extract-pptx.py <输入文件.pptx> [输出目录]

依赖：pip install python-pptx
"""

import json
import os
import sys
from pptx import Presentation


def extract_pptx(file_path, output_dir="."):
    """
    从 PowerPoint 文件中提取所有内容。
    返回包含文本、图片和备注的幻灯片数据字典列表。
    """
    prs = Presentation(file_path)
    slides_data = []

    # 创建用于存放提取图片的 assets 目录
    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            "number": slide_num + 1,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
        }

        for shape in slide.shapes:
            # 提取文本内容
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    slide_data["title"] = shape.text
                else:
                    slide_data["content"].append(
                        {"type": "text", "content": shape.text}
                    )

            # 提取图片
            if shape.shape_type == 13:  # 图片类型
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
                image_path = os.path.join(assets_dir, image_name)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                slide_data["images"].append(
                    {
                        "path": f"assets/{image_name}",
                        "width": shape.width,
                        "height": shape.height,
                    }
                )

        # 提取演讲者备注
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            slide_data["notes"] = notes_frame.text

        slides_data.append(slide_data)

    return slides_data


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法：python extract-pptx.py <输入文件.pptx> [输出目录]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."

    slides = extract_pptx(input_file, output_dir)

    # 将提取的数据写入 JSON 文件
    output_path = os.path.join(output_dir, "extracted-slides.json")
    with open(output_path, "w") as f:
        json.dump(slides, f, indent=2)

    print(f"已提取 {len(slides)} 张幻灯片到 {output_path}")
    for s in slides:
        img_count = len(s["images"])
        print(f"  幻灯片 {s['number']}：{s['title'] or '（无标题）'} — {img_count} 张图片")
